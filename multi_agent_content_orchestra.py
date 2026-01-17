# Core libraries and third-party imports used across the pipeline
from crewai import Crew, Task, Agent
from crewai_tools import SerperDevTool
from langchain_ibm import WatsonxLLM
from pptx import Presentation

import os

# Create folder for output files
OUTPUT_DIR = "content_output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ================================
# CONFIG
# ================================
# Recommended: set these in your terminal before running, e.g.:
# export WATSONX_APIKEY="your-real-key"
# export SERPER_API_KEY="your-real-serper-key"

# NOTE: Hardcoding secrets is not recommended in production.
# Prefer using environment variables in your shell or a secure secret manager.
os.environ["WATSONX_APIKEY"] = "YOUR_WATSONX_API_KEY_HERE"
os.environ["SERPER_API_KEY"] = "YOUR_SERPER_API_KEY_HERE"

# Common watsonx parameters used for all LLM calls
# - decoding_method="greedy" ensures deterministic output
# - max_new_tokens controls the maximum length of generated text
params = {"decoding_method": "greedy", "max_new_tokens": 500}

# IBM Cloud watsonx.ai config (Frankfurt / eu-de)
# URL and PROJECT_ID correspond to your IBM Cloud project and region.
#In the workspace where you create the watsonx_apikey, create a project and copy the project id here.
URL = "https://eu-de.ml.cloud.ibm.com"
PROJECT_ID = ""  # Amrita's sandbox (with associated service)


# Topic for the content pipeline – user-provided at runtime
TOPIC = input("Enter the topic you want to generate content for: ").strip()

# Safety fallback in case user hits enter
if not TOPIC:
    raise ValueError("Topic cannot be empty. Please run again and enter a valid topic.")


# ================================
# LLMs
# ================================

# Main reasoning / writing LLM
# This is the primary model used for research and core writing.
llm = WatsonxLLM(
    model_id="meta-llama/llama-3-3-70b-instruct",  # supported in your env
    url=URL,
    params=params,
    project_id=PROJECT_ID,
    verify=False,  # TODO: remove when SSL certs are configured properly (enables SSL verification)
)

# Function-calling / tools LLM (smaller, more structured)
# This model is used when tools (like web search) need to be invoked deterministically.
function_calling_llm = WatsonxLLM(
    model_id="mistralai/mistral-small-3-1-24b-instruct-2503",
    url=URL,
    params={**params, "temperature": 0.1},  # lower temperature -> more deterministic behavior
    project_id=PROJECT_ID,
    verify=False,
)

# Bigger token budget for formatter/editor
# This variant has a higher max_new_tokens limit to handle longer edits/formatting tasks.
fmt_llm = WatsonxLLM(
    model_id="meta-llama/llama-3-3-70b-instruct",
    url=URL,
    params={"decoding_method": "greedy", "max_new_tokens": 1500},
    project_id=PROJECT_ID,
    verify=False,
)


# ================================
# TOOLS
# ================================

# Web search tool (Serper API) used by the researcher agent to fetch live information.
#You may add more tools here as per the requirement.
search = SerperDevTool()


# ================================
# AGENTS
# ================================

# 1) Researcher – finds and structures real information
# This agent is responsible for doing web research and extracting key insights.
researcher = Agent(
    llm=llm,
    function_calling_llm=function_calling_llm,
    role="Senior AI Researcher",
    goal=f"Research the topic '{TOPIC}' using the web and extract the most important, recent, and practical insights.",
    backstory=(
        "You are a veteran AI researcher with experience in industry and academia. "
        "You care about credible sources, practical relevance, and clear structure."
    ),
    allow_delegation=False,  # this agent does not delegate work to other agents
    tools=[search],          # the web search tool is available to this agent
    verbose=1,               # verbosity controls how much logging CrewAI prints
)

# 2) Core Writer – builds the master narrative
# This agent takes the research notes and turns them into a cohesive explanation.
core_writer = Agent(
    llm=llm,
    role="Content Architect",
    goal=(
        f"Turn the research about '{TOPIC}' into a clear, structured explanation that can be reused in multiple formats."
    ),
    backstory=(
        "You are excellent at teaching complex technical concepts to mixed audiences. "
        "You write in a clear, logical, and engaging way."
    ),
    allow_delegation=False,
    verbose=1,
)

# 3) Formatter – creates the 6 content formats
# This agent repurposes the core narrative into multiple content types (speech, post, script, slides, etc.).
formatter = Agent(
    llm=fmt_llm,
    role="Multi-format Content Strategist",
    goal=(
        f"Convert the core narrative about '{TOPIC}' into multiple content formats "
        "that are immediately usable on stage, on LinkedIn, in video, and in business contexts."
    ),
    backstory=(
        "You are a content strategist who knows how to repurpose one strong idea into many formats. "
        "You understand how speakers, founders, and leaders communicate on different channels."
    ),
    allow_delegation=False,
    verbose=1,
)

# 4) Editor – polishes all outputs
# This agent is the final quality pass to refine tone, clarity, and consistency across all formats.
editor = Agent(
    llm=fmt_llm,
    role="Senior Editor",
    goal=(
        "Polish all generated content for clarity, conciseness, and impact while preserving meaning. "
        "Make sure everything flows well and feels coherent as one content pack."
    ),
    backstory=(
        "You are an experienced editor who has worked on talks, posts, decks, and executive summaries. "
        "You eliminate fluff and keep the strongest ideas and phrasing."
    ),
    allow_delegation=False,
    verbose=1,
)

# ================================
# TASKS
# ================================

# Task 1 – Research the topic using the web
# Input: TOPIC
# Output: structured research notes in bullet form, saved to 01_research_notes.txt
task_research = Task(
    description=(
        f"Use web search to research the topic '{TOPIC}'. "
        "Identify 5–10 key insights, trends, or important points.\n\n"
        "For each insight, include:\n"
        "- A short title\n"
        "- 2–3 sentence explanation\n"
        "- Why it matters in practice\n"
        "If possible, briefly reference where the insight comes from (e.g., company, paper, or blog)."
    ),
    expected_output=(
        "A structured set of 5–10 bullet points. Each bullet should include: "
        "a title, explanation, and why it matters. This will be used as the factual base."
    ),
    output_file=f"{OUTPUT_DIR}/01_research_notes.txt",  # Task output is written to this file
    agent=researcher,                      # Agent responsible for this task
)

# Task 2 – Create a core narrative (master explanation)
# Uses the research notes from Task 1 to produce an 800–1200 word narrative.
task_core_narrative = Task(
    description=(
        f"Using the research notes, write a clear, structured explanation of '{TOPIC}'. "
        "Assume the reader is technical but not an expert in agentic AI.\n\n"
        "Include sections:\n"
        "1. What is the topic and why now?\n"
        "2. Key concepts and frameworks (with simple analogies).\n"
        "3. Where it is used in practice (examples).\n"
        "4. Benefits and limitations.\n"
        "5. How this fits into the bigger AI landscape."
    ),
    expected_output=(
        "A 800–1200 word narrative that explains the topic clearly and can be reused for speech, posts, and scripts."
    ),
    output_file=f"{OUTPUT_DIR}/02_core_narrative.txt",
    agent=core_writer,
)

# Task 3 – Generate the multi-format content pack
# Transforms the core narrative into different content formats for various channels.
task_multiformat = Task(
    description=(
        "Based on the core narrative, create the following SIX distinct outputs.\n\n"
        "1) KEYNOTE SPEECH:\n"
        "- Structure: Introduction, 3–4 main sections, Conclusion\n"
        "- Tone: confident, inspiring, story-driven\n\n"
        "2) LINKEDIN POST:\n"
        "- Structure: 2-line hook, 4–6 lines of insight, 1-line takeaway or CTA\n"
        "- Tone: conversational, practical, shareable\n\n"
        "3) YOUTUBE SCRIPT (5–7 min):\n"
        "- Parts: Hook, Setup, Main Content (3 chapters), Recap, Outro/CTA\n"
        "- Include stage directions like [B-ROLL], [ON SCREEN TEXT] sparingly.\n\n"
        "4) SLIDE OUTLINE:\n"
        "- Provide slide titles and 3–5 bullets per slide.\n"
        "- Aim for 8–12 slides total.\n\n"
        "Clearly label each section with a heading like:\n"
        "[KEYNOTE SPEECH]\n"
        "[LINKEDIN POST]\n"
        "[YOUTUBE SCRIPT]\n"
        "[SLIDE OUTLINE]\n"
    ),
    expected_output=(
        "A single combined document containing all six sections, each clearly labeled with the specified headings."
    ),
    output_file= f"{OUTPUT_DIR}/03_multiformat_content_pack_raw.txt",
    agent=formatter,
)

# Task 4 – Edit and polish the full content pack
# The editor cleans up and refines the combined multi-format content.
task_edit = Task(
    description=(
        "Take the multi-format content pack and polish it for clarity, flow, and impact.\n\n"
        "Rules:\n"
        "- Keep the four sections and their headings.\n"
        "- Improve wording, remove repetition, and tighten long sentences.\n"
        "- Ensure the tone is consistent across formats (expert, friendly, clear).\n"
        "- Fix any obvious logical inconsistencies or contradictions."
    ),
    expected_output=(
        "A final, polished content pack with the same six labeled sections, ready for real-world use."
    ),
    output_file=f"{OUTPUT_DIR}/04_multiformat_content_pack_final.txt",
    agent=editor,
)

# ================================
# CREW
# ================================

# Crew orchestrates the agents and tasks in a pipeline.
# The order of tasks defines the flow of information.
crew = Crew(
    agents=[researcher, core_writer, formatter, editor],
    tasks=[task_research, task_core_narrative, task_multiformat, task_edit],
    verbose=1,  # higher verbosity -> more logging about task execution
)

def save_sections_to_files(full_text: str):
    """Split the final content pack into separate files per section.

    The function:
    - Looks for section markers (e.g., [KEYNOTE SPEECH])
    - Collects lines under each marker
    - Writes each section to its own .txt file
    """
    # Markers that identify the beginning of each section in the final content
    markers = [
        "[KEYNOTE SPEECH]",
        "[LINKEDIN POST]",
        "[YOUTUBE SCRIPT]",
        "[SLIDE OUTLINE]",
    ]

    # Dictionary mapping each marker to a list of lines
    sections = {m: [] for m in markers}
    current = None  # tracks which section we are currently populating

    # Iterate through the full text line by line
    for line in full_text.splitlines():
        stripped = line.strip()
        # If the line is exactly one of the markers, switch current section
        if stripped in markers:
            current = stripped
            continue
        # If we are within a known section, append the line to that section
        if current:
            sections[current].append(line)

    # Write each populated section to its own file
    for marker, lines in sections.items():
        if not lines:
            continue  # model might skip a section sometimes
        # Convert marker like "[KEYNOTE SPEECH]" -> "keynote_speech.txt"
        filename = marker.strip("[]").lower().replace(" ", "_") + ".txt"
        with open(f"{OUTPUT_DIR}/{filename}", "w") as f:
            f.write("\n".join(lines).strip())
        print(f"Saved {marker} to {filename}")


def build_slides_from_outline(outline_path=f"{OUTPUT_DIR}/slide_outline.txt", pptx_path=f"{OUTPUT_DIR}/slides.pptx"):
    """Build a PPTX deck from the slide outline text file.

    Expected outline format:
    - Each slide is a block separated by a blank line.
    - First line of block: 'Slide X: Title' or just 'Title'
    - Subsequent lines: bullet points (optionally prefixed with '-', '•', etc.)
    """
    # Check if the outline file exists before proceeding
    if not os.path.exists(outline_path):
        print(f"[slides] No slide outline file found at {outline_path}. Skipping PPTX generation.")
        return

    # Read the entire outline content
    with open(outline_path, "r") as f:
        text = f.read()

    # Split slides by blank lines into "blocks"
    blocks = [b.strip() for b in text.split("\n\n") if b.strip()]

    if not blocks:
        print("[slides] Slide outline is empty, nothing to build.")
        return

    # Initialize a new PowerPoint presentation
    prs = Presentation()

    # Process each slide block
    for block in blocks:
        # Split block into lines and clean whitespace
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue

        # First line: "Slide X: Title"  -> extract title after colon if present
        raw_title = lines[0]
        if ":" in raw_title:
            title = raw_title.split(":", 1)[1].strip()
        else:
            title = raw_title

        # Remaining lines are bullet points; strip any leading bullet characters
        bullet_lines = [l.lstrip("-• ").strip() for l in lines[1:]]

        # Use "Title and Content" layout (index 1) for each slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        # Add bullet points to the content placeholder
        body = slide.placeholders[1].text_frame
        first = True
        for bullet in bullet_lines:
            if not bullet:
                continue
            if first:
                # The first bullet sets the initial text of the text_frame
                body.text = bullet
                first = False
            else:
                # Subsequent bullets become new paragraphs
                p = body.add_paragraph()
                p.text = bullet

    # Save the generated PowerPoint file
    prs.save(pptx_path)
    print(f"[slides] ✅ PPTX generated: {pptx_path}")


if __name__ == "__main__":
    # Kick off the Crew pipeline:
    # 1) task_research
    # 2) task_core_narrative
    # 3) task_multiformat
    # 4) task_edit
    final_result = crew.kickoff()

    # Print the final consolidated result to stdout
    print("\n========== FINAL RESULT ==========\n")
    print(final_result)

    # Split the final content into separate text files by section markers
    save_sections_to_files(final_result)
    
    # Build slides.pptx from slide_outline.txt (if it exists and is well-formed)
    build_slides_from_outline(outline_path=f"{OUTPUT_DIR}/slide_outline.txt", pptx_path=f"{OUTPUT_DIR}/slides.pptx")
